# document_processor.py

import os
import fitz  # PyMuPDF
from pptx import Presentation
from langchain.text_splitter import RecursiveCharacterTextSplitter
from openai import OpenAI
from tqdm import tqdm
import requests
from requests.auth import HTTPBasicAuth
from bs4 import BeautifulSoup
from typing import Generator, List, Dict
from qdrant_client import QdrantClient
from qdrant_client.http import models
from PIL import Image
from io import BytesIO
from transformers import CLIPProcessor, CLIPModel


class FileProcessor:
    def process_pdf(self, file_path: str) -> Generator[tuple, None, None]:
        """Process a single PDF file using PyMuPDF."""
        try:
            document = fitz.open(file_path)
            for page_number in range(len(document)):
                page = document[page_number]
                text = page.get_text("text")
                if text.strip():
                    yield page_number, text.strip(), f"Header for page {page_number + 1}"
        except Exception as e:
            print(f"Error processing PDF {file_path}: {str(e)}")
            yield from []


    def process_pptx(self, file_path: str) -> Generator[tuple, None, None]:
        """Process a PPTX file."""
        try:
            presentation = Presentation(file_path)
            for slide_number, slide in enumerate(presentation.slides):
                slide_text = []
                slide_title = ""
                if slide.shapes.title:
                    slide_title = slide.shapes.title.text.strip()
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text = shape.text.strip()
                        if text and text != slide_title:
                            slide_text.append(text)
                main_text = "\n".join(slide_text)
                if main_text or slide_title:
                    yield slide_number, main_text, slide_title
        except Exception as e:
            print(f"Error processing PPTX {file_path}: {str(e)}")
            yield from []


    def process_document(self, file_path: str) -> Generator[tuple, None, None]:
        """Process any supported document type."""
        file_extension = os.path.splitext(file_path)[1].lower()
        if file_extension == '.pdf':
            yield from self.process_pdf(file_path)
        elif file_extension in ['.pptx', '.ppt']:
            yield from self.process_pptx(file_path)
        else:
            print(f"Unsupported file type: {file_extension}")
            yield from []


class ConfluenceProcessor:
    def __init__(self, user_name: str, api_token: str):
        self.api_token = api_token
        self.user_name = user_name


    def process_confluence(self, page_id: str) -> Generator[tuple, None, None]:
        """Fetch and process content from a Confluence page."""
        confluence_url = f'https://radhatrial.atlassian.net/wiki/rest/api/content/{page_id}?expand=body.storage'
        try:
            response = requests.get(confluence_url, auth=HTTPBasicAuth(self.user_name, self.api_token))
            if response.status_code == 200:
                data = response.json()
                page_content = data['body']['storage']['value']
                soup = BeautifulSoup(page_content, 'html.parser')
                text = soup.get_text()
                images = [img['src'] for img in soup.find_all('img')]
                yield 0, text.strip(), data['title'], images
            else:
                print(f"Failed to retrieve Confluence page: {response.status_code} - {response.text}")
        except Exception as e:
            print(f"Error processing Confluence page: {str(e)}")
            yield from []


class VectorStore:
    def __init__(self, config):
        self.qdrant_client = QdrantClient(path=config.LOCAL_QDRANT_PATH)
        self.collection_name = config.COLLECTION_NAME
        self._setup_collection()


    def _setup_collection(self):
        """Setup Qdrant collection if it doesn't exist."""
        collections = self.qdrant_client.get_collections()
        if not any(c.name == self.collection_name for c in collections.collections):
            self.qdrant_client.create_collection(
                collection_name=self.collection_name,
                vectors_config=models.VectorParams(
                    size=1536,
                    distance=models.Distance.COSINE
                )
            )


    def store_vectors(self, vectors: List[List[float]], chunks: List[Dict]):
        """Store vectors and metadata in Qdrant."""
        try:
            points = [
                models.PointStruct(
                    id=abs(hash(f"{chunk['metadata']['filename']}_{chunk['metadata']['page_number']}_{chunk['metadata']['chunk_number']}")) % (2**63),
                    vector=vector,
                    payload={
                        'text': chunk['text'],
                        'filename': chunk['metadata']['filename'],
                        'page_number': chunk['metadata']['page_number'],
                        'chunk_number': chunk['metadata']['chunk_number'],
                        'page_header': chunk['metadata']['page_header']
                    }
                )
                for chunk, vector in zip(chunks, vectors)
            ]
            self.qdrant_client.upsert(
                collection_name=self.collection_name,
                points=points
            )
        except Exception as e:
            print(f"Error storing vectors: {str(e)}")
            raise


class DocumentProcessor:
    def __init__(self, config):
        """Initialize document processor with all necessary components."""
        self.config = config
        self.openai_client = OpenAI(api_key=config.OPENAI_API_KEY)
        self.file_processor = FileProcessor()
        self.confluence_processor = ConfluenceProcessor(config.CONFLUENCE_USERNAME, config.CONFLUENCE_API_TOKEN)
        self.vector_store = VectorStore(config)

        # Text splitter for chunking
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=config.CHUNK_SIZE,
            chunk_overlap=config.CHUNK_OVERLAP,
            length_function=len,
            separators=["\n\n", "\n", ".", " ", ""]
        )


    def _get_processed_files(self) -> set:
        """Get set of already processed files."""
        try:
            response = self.vector_store.qdrant_client.scroll(
                collection_name=self.vector_store.collection_name,
                limit=10000,
                with_payload=['filename'],
                with_vectors=False
            )
            processed_files = {point.payload['filename'] for point in response[0]}
            processed_files.add(f"confluence_page_{self.config.CONFLUENCE_PAGE_ID}")  # Track Confluence page
            print("\n\nDebugging: processed_files : ", processed_files)
            return processed_files
        except Exception:
            return set()


    def create_chunks(self, text: str, metadata: Dict) -> List[Dict]:
        """Create chunks from text with metadata."""
        chunks = self.text_splitter.split_text(text)
        return [
            {
                'text': chunk,
                'metadata': {
                    'filename': metadata['filename'],
                    'page_number': metadata['page_number'],
                    'chunk_number': i + 1,
                    'page_header': metadata['page_header']
                }
            }
            for i, chunk in enumerate(chunks)
        ]


    def get_embeddings(self, texts: List[str]) -> List[List[float]]:
        """Get embeddings for texts using OpenAI."""
        try:
            response = self.openai_client.embeddings.create(
                model="text-embedding-ada-002",
                input=texts
            )
            return [data.embedding for data in response.data]
        except Exception as e:
            print(f"Error generating embeddings: {str(e)}")
            raise


    def get_image_embeddings(self, image_urls: List[str]) -> List[List[float]]:
        """Get embeddings for images using the CLIP model."""

        # Load the CLIP model and processor
        model = CLIPModel.from_pretrained("openai/clip-vit-base-patch16")
        processor = CLIPProcessor.from_pretrained("openai/clip-vit-base-patch16")

        embeddings = []

        for url in image_urls:
            try:
                # Retrieve the image
                response = requests.get(url)
                if response.status_code == 200:
                    image_data = response.content
                    
                    # Process the image
                    image = Image.open(BytesIO(image_data)).convert("RGB")  # Ensure image is in RGB format

                    # Generate embeddings
                    inputs = processor(images=image, return_tensors="pt")
                    with torch.no_grad():
                        image_embedding = model.get_image_features(**inputs)

                    # Convert the tensor to a list
                    image_embedding = image_embedding.squeeze().tolist()  # Convert to list and remove single-dimensional entries
                else:
                    print(f"Failed to retrieve image from {url}, status code: {response.status_code}")
                    image_embedding = [0] * 1536  # Dummy embedding for failed retrieval
                
            except Exception as e:
                print(f"Error processing image from {url}: {str(e)}")
                image_embedding = [0] * 1536  # Dummy embedding for errors

            embeddings.append(image_embedding)

        return embeddings


    def final_processing(self) -> int:
        """Main processing function."""
        processed_files = self._get_processed_files()

        # Get files from both PDF and PPT directories
        pdf_files = {f for f in os.listdir(self.config.PDF_DIRECTORY)
                     if f.endswith('.pdf')}
        ppt_files = {f for f in os.listdir(self.config.PPT_DIRECTORY)
                     if f.endswith(('.ppt', '.pptx'))}

        current_files = pdf_files.union(ppt_files)
        print("\n\nDebugging: current_files : ", current_files)

        new_files = current_files - processed_files
        print("\n\nDebugging: new_files : ", new_files)

        # if not new_files:
        #     print("No new files to process")
        #     return 0

        total_chunks = []
        total_processed = 0

        # Process each file
        for filename in tqdm(new_files, desc="Processing files"):
            if filename.endswith('.pdf'):
                file_path = os.path.join(self.config.PDF_DIRECTORY, filename)
            else:  # .ppt or .pptx
                file_path = os.path.join(self.config.PPT_DIRECTORY, filename)

            print(f"\nProcessing {filename}")

            # Process each page
            for page_number, text, page_header in self.file_processor.process_document(file_path):
                if text:
                    chunks = self.create_chunks(text, {
                        'filename': filename,
                        'page_number': page_number + 1,
                        'page_header': page_header
                    })
                    total_chunks.extend(chunks)

                # Process in batches
                if len(total_chunks) >= self.config.BATCH_SIZE:
                    try:
                        batch = total_chunks[:self.config.BATCH_SIZE]
                        embeddings = self.get_embeddings([c['text'] for c in batch])
                        self.vector_store.store_vectors(embeddings, batch)
                        total_processed += len(batch)
                        total_chunks = total_chunks[self.config.BATCH_SIZE:]
                        print(f"Processed {total_processed} chunks so far...")
                    except Exception as e:
                        print(f"Error processing batch: {str(e)}")
                        continue

        # Process remaining chunks
        if total_chunks:
            try:
                embeddings = self.get_embeddings([c['text'] for c in total_chunks])
                self.vector_store.store_vectors(embeddings, total_chunks)
                total_processed += len(total_chunks)
            except Exception as e:
                print(f"Error processing final batch: {str(e)}")

        # Process Confluence content
        confluence_processed = False
        for _, text, title, images in self.confluence_processor.process_confluence(self.config.CONFLUENCE_PAGE_ID):
            if text:
                chunks = self.create_chunks(text, {
                    'filename': title,
                    'page_number': 0,  # Use 0 for Confluence content
                    'page_header': title
                })
                embeddings = self.get_embeddings([c['text'] for c in chunks])
                self.vector_store.store_vectors(embeddings, chunks)
                confluence_processed = True  # Mark as processed

                        # Process images
            if images:
                image_embeddings = self.get_image_embeddings(images)
                for img_url, img_embedding in zip(images, image_embeddings):
                    self.vector_store.store_vectors([img_embedding], [{
                        'text': f"Image from {title}",
                        'metadata': {
                            'filename': title,
                            'page_number': 0,
                            'chunk_number': 0,  # Placeholder for image chunk number
                            'page_header': title
                        }
                    }])
                confluence_processed = True  # Mark as processed

        if confluence_processed:
            print(f"Processed Confluence page {self.config.CONFLUENCE_PAGE_ID}.")

        return total_processed


def main():
    from config import Config

    # Initialize directory structure
    directories = [
        Config.DOCUMENT_DIRECTORY,
        Config.PDF_DIRECTORY,
        Config.PPT_DIRECTORY,
        Config.LOCAL_QDRANT_PATH
    ]

    for directory in directories:
        os.makedirs(directory, exist_ok=True)

    processor = DocumentProcessor(Config)
    try:
        print("Starting document processing...")
        num_processed = processor.final_processing()
        print(f"Successfully processed {num_processed} chunks")
    except Exception as e:
        print(f"Error processing documents: {str(e)}")


if __name__ == "__main__":
    main()